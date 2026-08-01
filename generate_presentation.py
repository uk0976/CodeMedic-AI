import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Define Color Palette (Dark Obsidian Glassmorphism)
BG_DARK = RGBColor(11, 15, 25)           # #0B0F19 Deep Obsidian Background
CARD_BG = RGBColor(19, 28, 49)          # #131C31 Dark Slate Card Fill
CARD_BG_ALT = RGBColor(15, 23, 42)      # #0F172A Dark Slate Card Alt
PRIMARY_BLUE = RGBColor(37, 99, 235)    # #2563EB Vibrant Blue
CYAN = RGBColor(6, 182, 212)            # #06B6D4 Electric Cyan
SKY_BLUE = RGBColor(56, 189, 248)       # #38BDF8 Sky Blue
PURPLE = RGBColor(124, 58, 237)         # #7C3AED Neon Purple
PURPLE_DARK = RGBColor(46, 16, 101)     # #2E1065 Deep Purple Fill
INDIGO = RGBColor(79, 70, 229)          # #4F46E5 Indigo Accent
AMBER = RGBColor(245, 158, 11)          # #F59E0B Amber Accent
GREEN = RGBColor(52, 211, 153)          # #34D399 Mint Green
RED = RGBColor(248, 113, 113)           # #F87171 Coral Red

TEXT_WHITE = RGBColor(248, 250, 252)    # #F8FAFC Crisp White Text
TEXT_MUTED = RGBColor(148, 163, 184)    # #94A3B8 Slate 400 Muted Text
TEXT_DIM = RGBColor(100, 116, 139)      # #64748B Dimmed Text
BORDER_GLOW = RGBColor(30, 58, 138)     # #1E3A8A Dark Blue Glow Border
BORDER_PURPLE = RGBColor(109, 40, 217)  # #6D28D9 Neon Purple Border
BORDER_CYAN = RGBColor(14, 116, 144)    # #0E7490 Cyan Border

# Font Family: Times New Roman
FONT_TITLE = "Times New Roman"
FONT_BODY = "Times New Roman"

MEDIA_DIR = os.path.join(os.path.dirname(__file__), "extracted_media")

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_header(slide, badge_text, title_text, subtitle_text=""):
        # Category Badge
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.35))
        tf = badge_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = badge_text.upper()
        p.font.name = FONT_TITLE
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = CYAN

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(9.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.32), Inches(10.0), Inches(0.4))
            tf = sub_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.italic = True
            p.font.color.rgb = TEXT_MUTED

    def add_anim_badge(slide, text):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(0.4), Inches(3.1), Inches(0.38))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_PURPLE
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✨ ANIMATION: {text}"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = SKY_BLUE
        p.alignment = PP_ALIGN.CENTER

    def add_glass_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_GLOW):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    # ==========================================
    # SLIDE 1: Hero Cover
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)
    add_anim_badge(slide1, "Hero Fade & Gradient Glow")

    # Main Hero Glass Card
    hero_card = add_glass_card(slide1, Inches(0.8), Inches(1.0), Inches(11.733), Inches(5.6), bg_color=CARD_BG, border_color=BORDER_GLOW)
    
    # Left decorative accent bar
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(0.2), Inches(5.6))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_BLUE
    accent_bar.line.fill.background()

    # Track Badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.4), Inches(4.7), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = CARD_BG_ALT
    badge.line.color.rgb = BORDER_CYAN
    badge.line.width = Pt(1)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "OPENAI CODEX HACKATHON 2026 • TRACK: AGENTIC CODING"
    p.font.name = FONT_TITLE
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(7.5), Inches(1.1))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CodeMedic AI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE

    # Tagline
    tag_box = slide1.shapes.add_textbox(Inches(1.3), Inches(3.1), Inches(7.5), Inches(0.6))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Fix. Explain. Optimize."
    p.font.name = FONT_TITLE
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # Powered By
    p_box = slide1.shapes.add_textbox(Inches(1.3), Inches(3.7), Inches(7.5), Inches(0.5))
    tf = p_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ Powered by OpenAI Codex"
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Description Paragraph
    desc_box = slide1.shapes.add_textbox(Inches(1.3), Inches(4.3), Inches(7.0), Inches(1.1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "An autonomous multi-agent AI assistant that revolutionizes software engineering by combining real-time bug detection, security scanning, performance optimization, and automated test & doc generation into a unified workflow."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # Developer Info Card (Right Side)
    dev_card = add_glass_card(slide1, Inches(8.7), Inches(1.8), Inches(3.4), Inches(4.0), bg_color=CARD_BG_ALT, border_color=BORDER_PURPLE)
    dev_tf = dev_card.text_frame
    dev_tf.word_wrap = True
    
    p = dev_tf.paragraphs[0]
    p.text = "DEVELOPER"
    p.font.name = FONT_TITLE
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    
    p2 = dev_tf.add_paragraph()
    p2.text = "Umer Khan"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    
    p3 = dev_tf.add_paragraph()
    p3.text = "B.E. Artificial Intelligence & Data Science"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED

    p4 = dev_tf.add_paragraph()
    p4.text = "\nEVENT & TRACK"
    p4.font.name = FONT_TITLE
    p4.font.size = Pt(10.5)
    p4.font.bold = True
    p4.font.color.rgb = PURPLE

    p5 = dev_tf.add_paragraph()
    p5.text = "OpenAI Codex Hackathon 2026\nTrack: Agentic Coding"
    p5.font.name = FONT_BODY
    p5.font.size = Pt(11.5)
    p5.font.bold = True
    p5.font.color.rgb = CYAN

    # ==========================================
    # SLIDE 2: Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "PROBLEM STATEMENT", "Current Challenges in Software Engineering", "Software developers face severe productivity bottlenecks with fragmented manual tools.")
    add_anim_badge(slide2, "Staggered Card Entrance & Icon Pop")

    problems = [
        ("⏱️", "Time Consuming Debugging", "Manual bug hunting and root-cause analysis consume up to 50% of active development cycles."),
        ("🔄", "Context Switching", "Juggling static analysis tools, review platforms, and documentation causes severe focus loss."),
        ("🛡️", "Security Risks", "Critical security vulnerabilities remain undetected until staging or production releases."),
        ("📚", "Manual Documentation", "Code explanations and documentation are frequently incomplete, outdated, or missing entirely."),
        ("⏳", "Delayed Reviews", "Heavy reliance on senior developers for code reviews creates release schedule bottlenecks."),
        ("⚡", "Performance Bottlenecks", "Algorithmic inefficiencies and memory bottlenecks are hard to spot without deep profiling."),
        ("📉", "Productivity Loss", "Repetitive manual engineering overhead directly slows down feature delivery cycles.")
    ]

    card_w = Inches(2.75)
    card_h = Inches(2.2)
    gap_x = Inches(0.2)
    start_y = Inches(1.9)

    for i, (icon, title, desc) in enumerate(problems[:4]):
        left = Inches(0.8) + i * (card_w + gap_x)
        card = add_glass_card(slide2, left, start_y, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED

    start_y2 = Inches(4.3)
    for i, (icon, title, desc) in enumerate(problems[4:]):
        left = Inches(0.8) + i * (card_w + gap_x)
        card = add_glass_card(slide2, left, start_y2, card_w, card_h, bg_color=CARD_BG, border_color=BORDER_PURPLE)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED

    # Highlight Summary Card (Last column bottom)
    sum_left = Inches(0.8) + 3 * (card_w + gap_x)
    sum_card = add_glass_card(slide2, sum_left, start_y2, card_w, card_h, bg_color=CARD_BG_ALT, border_color=CYAN)
    stf = sum_card.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = Inches(0.15)
    p = stf.paragraphs[0]
    p.text = "🎯 CORE IMPACT"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p2 = stf.add_paragraph()
    p2.text = "\nConventional development flows depend on fragmented, manual tools—resulting in high bug leakage and slow delivery."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10.5)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 3: Existing Workflow
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, "EXISTING WORKFLOW", "Fragmented & Manual Development Process", "The traditional workflow relies on disconnected tools and slow manual review cycles.")
    add_anim_badge(slide3, "Sequential Workflow Step Highlight")

    # Add PDF Image 2
    img2_path = os.path.join(MEDIA_DIR, "image2.png")
    if os.path.exists(img2_path):
        slide3.shapes.add_picture(img2_path, Inches(0.8), Inches(1.9), width=Inches(7.2))

    # Right side 7-step callout card
    right_card = add_glass_card(slide3, Inches(8.3), Inches(1.9), Inches(4.2), Inches(5.1), bg_color=CARD_BG, border_color=BORDER_GLOW)
    rtf = right_card.text_frame
    rtf.word_wrap = True
    rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = Inches(0.2)

    p = rtf.paragraphs[0]
    p.text = "TRADITIONAL WORKFLOW STEPS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE

    steps = [
        "1. Write & Compile Code: Manual development in IDE.",
        "2. Manual Debugging: Line-by-line problem hunting.",
        "3. Web Search: Searching documentation for solutions.",
        "4. Code Review: Waiting for senior dev review feedback.",
        "5. Security & Perf: Running separate compliance scanners.",
        "6. Manual Tests & Docs: Writing boilerplate unit tests.",
        "7. Final Deployment: High risk of undetected bugs."
    ]

    for step in steps:
        p_step = rtf.add_paragraph()
        p_step.text = f"• {step}"
        p_step.font.name = FONT_BODY
        p_step.font.size = Pt(10.5)
        p_step.font.color.rgb = TEXT_MUTED

    # Bottom summary callout
    p_sum = rtf.add_paragraph()
    p_sum.text = "\n⚠️ Limitation: High context switching, long feedback loops, and delayed releases."
    p_sum.font.name = FONT_BODY
    p_sum.font.size = Pt(10.5)
    p_sum.font.bold = True
    p_sum.font.color.rgb = RED

    # ==========================================
    # SLIDE 4: Existing Solutions & Limitations
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, "SOLUTION LIMITATIONS", "Existing Solutions & Their Key Limitations", "Traditional tools offer rigid static rules or generic chatbot responses without structured context.")
    add_anim_badge(slide4, "Morph & Limitation Card Slide-In")

    tools = [
        ("ESLint / Linters", "Rule-Based Only", "Strictly checks static syntax rules without reasoning about complex programming logic, context, or design flaws.", RED),
        ("SonarQube", "Limited Context", "Detects basic code smells but lacks deep AI reasoning to suggest intelligent, context-aware architectural fixes.", AMBER),
        ("IDE Warnings", "Syntax-Focused", "Flags immediate syntax errors and missing imports, but fails to identify security risks or performance bottlenecks.", AMBER),
        ("Generic AI Chatbots", "Unstructured Output", "Requires manual prompt engineering, lacks multi-file context, and does not provide unified technical engineering reports.", RED),
        ("Manual Review", "Slow & Inconsistent", "Highly dependent on senior developer bandwidth, causing release delays, reviewer fatigue, and human error.", RED)
    ]

    c_w = Inches(2.2)
    c_h = Inches(4.8)
    gap = Inches(0.18)

    for i, (name, tag, limit, color) in enumerate(tools):
        left = Inches(0.8) + i * (c_w + gap)
        card = add_glass_card(slide4, left, Inches(1.9), c_w, c_h, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_TITLE
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        # Badge
        p_b = tf.add_paragraph()
        p_b.text = f"\n[{tag.upper()}]"
        p_b.font.name = FONT_TITLE
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{limit}"
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_MUTED

    # Bottom summary box
    b_box = add_glass_card(slide4, Inches(0.8), Inches(6.85), Inches(11.733), Inches(0.45), bg_color=CARD_BG_ALT, border_color=CYAN)
    btf = b_box.text_frame
    p = btf.paragraphs[0]
    p.text = "💡 CodeMedic AI Bridge: Replaces fragmented tools with an autonomous, multi-agent AI engine that provides end-to-end code analysis."
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 5: Our Solution
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, "PROPOSED SOLUTION", "CodeMedic AI — Agentic AI Engineering Assistant", "An intelligent multi-agent platform that automates code analysis, security, optimization, and testing.")
    add_anim_badge(slide5, "Agent Flow Pulse & Callout Reveal")

    # Add PDF Image 3
    img3_path = os.path.join(MEDIA_DIR, "image3.png")
    if os.path.exists(img3_path):
        slide5.shapes.add_picture(img3_path, Inches(0.8), Inches(1.9), width=Inches(7.2))

    # Right side 5 Pillars Card
    sol_card = add_glass_card(slide5, Inches(8.3), Inches(1.9), Inches(4.2), Inches(5.1), bg_color=CARD_BG, border_color=BORDER_PURPLE)
    stf = sol_card.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = Inches(0.2)

    p = stf.paragraphs[0]
    p.text = "WORKFLOW & KEY PILLARS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    flow_p = stf.add_paragraph()
    flow_p.text = "Developer ➔ Upload Code ➔ Agentic AI Engine ➔ Smart Dashboard\n"
    flow_p.font.name = FONT_BODY
    flow_p.font.size = Pt(10.5)
    flow_p.font.bold = True
    flow_p.font.color.rgb = CYAN

    pillars = [
        ("🐛 Bug Detection", "Identifies syntax, logical, and runtime edge cases."),
        ("🛡️ Security Vulnerabilities", "Scans for SQLi, secrets, and OWASP risks."),
        ("⚡ Performance Optimization", "Evaluates complexity and memory bottlenecks."),
        ("📚 AI Documentation", "Generates instant code explanations & comments."),
        ("🧪 Unit Test Generation", "Creates complete test suites automatically.")
    ]

    for title, desc in pillars:
        p_pil = stf.add_paragraph()
        p_pil.text = f"• {title}: {desc}"
        p_pil.font.name = FONT_BODY
        p_pil.font.size = Pt(10)
        p_pil.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 6: Hackathon Track Alignment
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "HACKATHON TRACK ALIGNMENT", "Agentic Coding — Perfect Alignment", "Demonstrating how autonomous AI agents collaborate to solve complex engineering challenges.")
    add_anim_badge(slide6, "Multi-Agent Network Graph Expansion")

    # Top Architecture Banner
    top_banner = add_glass_card(slide6, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.6), bg_color=CARD_BG_ALT, border_color=BORDER_PURPLE)
    ttf = top_banner.text_frame
    p = ttf.paragraphs[0]
    p.text = "🤖 Orchestrator Agent   ➔   👥 Multiple Specialized AI Agents   ➔   📊 Unified Engineering Report"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # 6 Pillar Alignment Cards (3x2 Grid)
    align_items = [
        ("🤖 Autonomous Code Analysis", "AI agents independently inspect submitted code without requiring manual prompting or step-by-step guidance."),
        ("🧩 Task Decomposition", "Complex software engineering tasks are automatically divided among specialized domain-specific agents."),
        ("🤝 Multi-Agent Collaboration", "Multiple AI agents work concurrently under Orchestrator guidance to analyze different software aspects."),
        ("🧠 Intelligent Decision Making", "The system dynamically determines the best analysis strategy based on uploaded code structure and language."),
        ("📊 Structured Engineering Reports", "Outputs from all agents are synthesized into a single, comprehensive, actionable developer report."),
        ("🚀 Developer Productivity", "Automates repetitive tasks—debugging, security, performance, docs, and testing—accelerating development velocity.")
    ]

    cw = Inches(3.75)
    ch = Inches(2.1)
    gx = Inches(0.24)
    gy = Inches(0.25)

    for i, (title, desc) in enumerate(align_items):
        row = i // 3
        col = i % 3
        left = Inches(0.8) + col * (cw + gx)
        top = Inches(2.65) + row * (ch + gy)

        card = add_glass_card(slide6, left, top, cw, ch, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(12.5)
        p.font.bold = True
        p.font.color.rgb = SKY_BLUE

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 7: System Architecture
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "SYSTEM ARCHITECTURE", "End-to-End Multi-Agent Data Flow", "A modular multi-agent pipeline designed for high-performance code analysis.")
    add_anim_badge(slide7, "Data Packet Animation Across Architecture")

    # Add PDF Image 4
    img4_path = os.path.join(MEDIA_DIR, "image4.png")
    if os.path.exists(img4_path):
        slide7.shapes.add_picture(img4_path, Inches(0.8), Inches(1.9), width=Inches(7.4))

    # Right side Architecture Breakdown Cards
    arch_card = add_glass_card(slide7, Inches(8.5), Inches(1.9), Inches(4.0), Inches(5.1), bg_color=CARD_BG, border_color=BORDER_GLOW)
    atf = arch_card.text_frame
    atf.word_wrap = True
    atf.margin_left = atf.margin_top = atf.margin_right = atf.margin_bottom = Inches(0.18)

    p = atf.paragraphs[0]
    p.text = "ARCHITECTURE LAYERS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = CYAN

    layers = [
        ("💻 Frontend Layer", "Next.js, TypeScript, Tailwind CSS, & Monaco Code Editor."),
        ("⚙️ Backend API", "FastAPI (Python) REST API for high-speed code processing."),
        ("🧠 Orchestrator Agent", "Parses code context and routes tasks to specialized agents."),
        ("🤖 Specialized Agents", "Parallel Bug, Security, Perf, Review, Test, & Doc agents."),
        ("📑 Synthesis Engine", "Aggregates multi-agent insights into a unified report schema."),
        ("📊 Smart Dashboard", "Displays health scores, issues, & downloadable PDF/MD/JSON.")
    ]

    for title, desc in layers:
        p_layer = atf.add_paragraph()
        p_layer.text = f"\n• {title}\n  {desc}"
        p_layer.font.name = FONT_BODY
        p_layer.font.size = Pt(10)
        p_layer.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 8: User Workflow
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, "USER WORKFLOW", "9-Step Seamless Analysis & Optimization Flow", "From login to report export, CodeMedic AI delivers an intuitive, developer-centric workflow.")
    add_anim_badge(slide8, "9-Step Sequential Flow Highlight")

    # Add PDF Image 5
    img5_path = os.path.join(MEDIA_DIR, "image5.png")
    if os.path.exists(img5_path):
        slide8.shapes.add_picture(img5_path, Inches(0.8), Inches(1.9), width=Inches(7.2))

    # Right side 9 Step Cards
    step_card = add_glass_card(slide8, Inches(8.3), Inches(1.9), Inches(4.2), Inches(5.1), bg_color=CARD_BG, border_color=BORDER_PURPLE)
    stf8 = step_card.text_frame
    stf8.word_wrap = True
    stf8.margin_left = stf8.margin_top = stf8.margin_right = stf8.margin_bottom = Inches(0.18)

    p = stf8.paragraphs[0]
    p.text = "EXECUTION WORKFLOW STEPS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    workflow_steps = [
        "1. Sign Up / Login: Access secure portal.",
        "2. Upload/Paste Code: Input code in Monaco Editor.",
        "3. AI Code Analysis: Orchestrator routes codebase.",
        "4. Multi-Agent Processing: Parallel agent execution.",
        "5. Results Aggregation: Synthesis Agent merges outputs.",
        "6. Smart Dashboard: View code health score & risks.",
        "7. Export Report: One-click PDF / MD / JSON export.",
        "8. Improve Code: Apply recommended fixes.",
        "9. Reanalyze: Perform continuous quality loop."
    ]

    for s in workflow_steps:
        ps = stf8.add_paragraph()
        ps.text = f"• {s}"
        ps.font.name = FONT_BODY
        ps.font.size = Pt(9.8)
        ps.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 9: Technical Stack
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    add_header(slide9, "TECHNICAL STACK", "Cutting-Edge Engineering & AI Ecosystem", "Built with modern frameworks, high-performance APIs, and OpenAI Codex.")
    add_anim_badge(slide9, "Stack Category Zoom & Glow")

    stack_categories = [
        ("🖥️ Frontend", "Next.js (React)\nTypeScript\nTailwind CSS\nshadcn/ui\nMonaco Editor\nLucide Icons"),
        ("⚙️ Backend", "FastAPI (Python)\nPython 3.12 Core\nREST API Services\nAsyncIO Processing\nPydantic Schemas"),
        ("🗄️ Database & Auth", "PostgreSQL\nSupabase Storage\nJWT Authentication\nSecure Session Tokens"),
        ("🧠 AI & LLM Engine", "OpenAI Codex\nChatGPT (GPT-4o)\nGroq API\nMulti-Agent Orchestrator"),
        ("🚀 Deployment", "Docker Containers\nVercel (Frontend)\nRender / Railway (Backend)\nGit & GitHub"),
        ("🛠️ Developer Tools", "VS Code IDE\nReportLab Engine\nMarkdown PDF Generator\nPostman API Tools")
    ]

    cw9 = Inches(3.75)
    ch9 = Inches(2.0)
    gx9 = Inches(0.24)
    gy9 = Inches(0.2)

    for i, (cat_title, cat_text) in enumerate(stack_categories):
        row = i // 3
        col = i % 3
        left = Inches(0.8) + col * (cw9 + gx9)
        top = Inches(1.9) + row * (ch9 + gy9)

        card = add_glass_card(slide9, left, top, cw9, ch9, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.12)

        p = tf.paragraphs[0]
        p.text = cat_title
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SKY_BLUE

        p_desc = tf.add_paragraph()
        p_desc.text = cat_text
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = TEXT_MUTED

    # Bottom Special Hackathon Banner
    h_banner = add_glass_card(slide9, Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.9), bg_color=CARD_BG_ALT, border_color=BORDER_PURPLE)
    htf = h_banner.text_frame
    htf.word_wrap = True
    htf.margin_left = htf.margin_top = htf.margin_right = htf.margin_bottom = Inches(0.1)

    p1 = htf.paragraphs[0]
    p1.text = "⚡ AI DEVELOPMENT & HACKATHON TOOLS: OpenAI Codex  |  ChatGPT  |  Antigravity"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = PURPLE

    p2 = htf.add_paragraph()
    p2.text = "Leveraged during the Vibe Coding Hackathon for AI-assisted coding, rapid prototyping, prompt engineering, debugging, documentation generation, and UI component creation."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 10: Core Features
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10)
    add_header(slide10, "CORE FEATURES", "12 Powerful Capabilities Driving Developer Productivity", "Completely automating code reviews, bug hunting, testing, and technical documentation.")
    add_anim_badge(slide10, "12-Card Cascade Reveal")

    # Add PDF Image 6 (Core Features visual banner)
    img6_path = os.path.join(MEDIA_DIR, "image6.png")
    if os.path.exists(img6_path):
        slide10.shapes.add_picture(img6_path, Inches(3.4), Inches(1.7), width=Inches(6.5), height=Inches(2.2))

    # 12 Feature Cards (4 columns x 3 rows) below/around
    features_12 = [
        ("🤖 1. Agentic AI Code Analysis", SKY_BLUE),
        ("🐛 2. Intelligent Bug Detection", RED),
        ("🛡️ 3. Security Vulnerability Analysis", AMBER),
        ("⚡ 4. Performance Optimization", GREEN),
        ("🔍 5. AI-Powered Code Review", PURPLE),
        ("🧪 6. Automatic Unit Test Gen", INDIGO),
        ("📚 7. AI Docs & Explanation", SKY_BLUE),
        ("📊 8. Interactive Smart Dashboard", CYAN),
        ("📑 9. Exportable PDF/MD Reports", PURPLE),
        ("🌐 10. Multi-Language Support", SKY_BLUE),
        ("💻 11. Professional Monaco Editor", TEXT_WHITE),
        ("🚀 12. Productivity Enhancement", GREEN)
    ]

    cw10 = Inches(2.75)
    ch10 = Inches(0.85)
    gx10 = Inches(0.24)
    start_y10 = Inches(4.1)

    for i, (feat_name, col_color) in enumerate(features_12):
        row = i // 4
        col = i % 4
        left = Inches(0.8) + col * (cw10 + gx10)
        top = start_y10 + row * (ch10 + Inches(0.15))

        card = add_glass_card(slide10, left, top, cw10, ch10, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = feat_name
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col_color
        p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 11: Innovation
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11)
    add_header(slide11, "KEY INNOVATIONS", "Breakthrough Capabilities in AI Software Engineering", "Pioneering the transition from single-prompt chatbots to collaborative Multi-Agent coding.")
    add_anim_badge(slide11, "Glass Card Morph & Pulse")

    innovations = [
        ("🤖 Agentic Coding Architecture", "Replaces single static AI models with a collaborative multi-agent network where specialized agents work together concurrently."),
        ("🎯 Orchestrator-Based Workflow", "A central AI Orchestrator intelligently parses the codebase, assigns tasks to domain agents, and synthesizes results."),
        ("🌐 All-in-One Development Platform", "Unifies debugging, security scanning, performance tuning, test generation, and documentation into a single dashboard."),
        ("📊 Interactive Smart Dashboard", "Presents complex AI analysis through an intuitive visual dashboard with health scores and actionable suggestions."),
        ("⚡ Automated Software Engineering", "Drastically cuts down manual engineering overhead, enabling developers to release high-quality code faster."),
        ("📄 One-Click Technical Reports", "Generates comprehensive, downloadable engineering documentation in PDF, Markdown, and JSON formats instantly."),
        ("🚀 Productivity Acceleration", "Transforms hours of manual code reviews and bug fixing into an automated, explainable AI workflow.")
    ]

    cw11 = Inches(3.75)
    ch11 = Inches(2.2)
    gx11 = Inches(0.24)
    gy11 = Inches(0.2)

    for i, (title, desc) in enumerate(innovations[:6]):
        row = i // 3
        col = i % 3
        left = Inches(0.8) + col * (cw11 + gx11)
        top = Inches(1.9) + row * (ch11 + gy11)

        card = add_glass_card(slide11, left, top, cw11, ch11, bg_color=CARD_BG, border_color=BORDER_PURPLE)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PURPLE

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED

    # Bottom Full-width 7th innovation card
    inn7 = add_glass_card(slide11, Inches(0.8), Inches(6.6), Inches(11.733), Inches(0.65), bg_color=CARD_BG_ALT, border_color=CYAN)
    itf = inn7.text_frame
    p = itf.paragraphs[0]
    p.text = f"💡 {innovations[6][0]}: {innovations[6][1]}"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 12: Comparison
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_bg(slide12)
    add_header(slide12, "SOLUTION COMPARISON", "CodeMedic AI vs Existing Tools & AI Chatbots", "Highlighting how CodeMedic AI outperforms traditional static analyzers and simple chatbots.")
    add_anim_badge(slide12, "Comparison Row-by-Row Highlight")

    # Table A: Feature Matrix
    t_shape = slide12.shapes.add_table(8, 4, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.2))
    table = t_shape.table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(1.3)
    table.columns[2].width = Inches(1.2)
    table.columns[3].width = Inches(1.4)

    headers = ["Feature", "Static Tools", "AI Chatbots", "CodeMedic AI"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = FONT_TITLE
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

    rows_data = [
        ("AI-Powered Analysis", "❌ Rule-only", "✅ Chat-based", "✅ Autonomous"),
        ("Multi-Agent Architecture", "❌ None", "❌ Single model", "✅ Multi-Agent"),
        ("Security Vulnerabilities", "✅ Static rules", "⚠️ Prompt-based", "✅ Automated SAST"),
        ("Unit Test Generation", "❌ Manual", "✅ Prompt-based", "✅ Automatic"),
        ("AI Code Documentation", "❌ None", "✅ Prompt-based", "✅ Auto-generated"),
        ("Interactive Dashboard", "❌ Static list", "❌ Chat log", "✅ Smart Dashboard"),
        ("Productivity Enhancement", "Moderate", "Moderate", "🚀 High Velocity")
    ]

    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else CARD_BG_ALT
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.bold = (j == 3 or j == 0)
            p.font.color.rgb = SKY_BLUE if j == 3 else TEXT_WHITE
            if j > 0:
                p.alignment = PP_ALIGN.CENTER

    # Table B: Workflow Shift (Right Side)
    card_right = add_glass_card(slide12, Inches(7.5), Inches(1.8), Inches(5.033), Inches(5.2), bg_color=CARD_BG, border_color=BORDER_PURPLE)
    crtf = card_right.text_frame
    crtf.word_wrap = True
    crtf.margin_left = crtf.margin_top = crtf.margin_right = crtf.margin_bottom = Inches(0.2)

    p = crtf.paragraphs[0]
    p.text = "TRADITIONAL WORKFLOW vs CODEMEDIC AI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    shifts = [
        ("Disconnected Tools", "Single Integrated Platform"),
        ("Manual Debugging", "AI-Assisted Bug Detection"),
        ("Separate Security Scanners", "Built-in Security Analysis"),
        ("Manual Perf Optimization", "AI Performance Tuning"),
        ("Manual Code Reviews", "Automated AI Code Review"),
        ("Manual Documentation", "Auto AI Doc Generation"),
        ("Manual Test Suite Creation", "Automatic Unit Test Gen"),
        ("Time-Consuming Workflow", "Accelerated AI Development")
    ]

    for old_w, new_w in shifts:
        ps = crtf.add_paragraph()
        ps.text = f"• {old_w}  ➔  {new_w}"
        ps.font.name = FONT_BODY
        ps.font.size = Pt(9.5)
        ps.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 13: Future Scope
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_bg(slide13)
    add_header(slide13, "FUTURE SCOPE", "Product Roadmap & Strategic Expansion", "Scaling CodeMedic AI to enterprise teams and native IDE ecosystems.")
    add_anim_badge(slide13, "Roadmap Node Travel & Pulse")

    roadmap_items = [
        ("🔌 VS Code Extension", "Perform instant multi-agent code analysis directly inside developer IDEs."),
        ("🐙 GitHub / GitLab Integration", "Automate pull request reviews and commit analysis via webhooks."),
        ("📁 Repository-Level Analysis", "Expand beyond single files to analyze entire multi-repository projects."),
        ("🔄 CI/CD Pipeline Integration", "Integrate into GitHub Actions, Azure DevOps, and Jenkins for build gates."),
        ("👥 Enterprise Dashboard", "Enable team-based reviews, shared analytics, and role-based access controls."),
        ("🧠 Self-Learning AI Agents", "Continuously refine model recommendations based on developer feedback loops."),
        ("📱 Mobile Companion App", "Review critical security alerts and analysis summaries on mobile devices.")
    ]

    cw13 = Inches(3.75)
    ch13 = Inches(2.1)
    gx13 = Inches(0.24)
    gy13 = Inches(0.2)

    for i, (title, desc) in enumerate(roadmap_items[:6]):
        row = i // 3
        col = i % 3
        left = Inches(0.8) + col * (cw13 + gx13)
        top = Inches(1.9) + row * (ch13 + gy13)

        card = add_glass_card(slide13, left, top, cw13, ch13, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SKY_BLUE

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED

    # Bottom 7th Roadmap Card
    rm7 = add_glass_card(slide13, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.65), bg_color=CARD_BG_ALT, border_color=BORDER_PURPLE)
    rtf7 = rm7.text_frame
    p = rtf7.paragraphs[0]
    p.text = f"📱 {roadmap_items[6][0]}: {roadmap_items[6][1]}"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 14: Demo Flow
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_bg(slide14)
    add_header(slide14, "PRODUCT DEMO", "Live Hackathon Demo Flow", "A 6-stage live demonstration walkthrough of CodeMedic AI in action.")
    add_anim_badge(slide14, "Demo Stage Transition Arrows")

    # Flow Bar
    flow_steps = [
        ("1. Landing Page", "Overview & Features", SKY_BLUE),
        ("2. Secure Login", "JWT Auth Portal", PURPLE),
        ("3. Monaco Editor", "Paste/Upload Code", INDIGO),
        ("4. AI Analysis", "Multi-Agent Engine", CYAN),
        ("5. Smart Dashboard", "Health Scores & Risks", GREEN),
        ("6. Export Report", "Download PDF / MD", AMBER)
    ]

    cw14 = Inches(1.8)
    gap14 = Inches(0.18)
    start_x14 = Inches(0.8)

    for i, (stage, desc, col_c) in enumerate(flow_steps):
        left = start_x14 + i * (cw14 + gap14)
        card = add_glass_card(slide14, left, Inches(1.9), cw14, Inches(4.8), bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.12)

        p = tf.paragraphs[0]
        p.text = f"STAGE 0{i+1}"
        p.font.name = FONT_TITLE
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = col_c

        p_t = tf.add_paragraph()
        p_t.text = f"\n{stage}"
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        p_d = tf.add_paragraph()
        p_d.text = f"\n{desc}"
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

        if i < 5:
            arrow = slide14.shapes.add_textbox(left + cw14, Inches(3.8), gap14, Inches(0.5))
            atf = arrow.text_frame
            ap = atf.paragraphs[0]
            ap.text = "➔"
            ap.font.size = Pt(14)
            ap.font.color.rgb = SKY_BLUE

    # Bottom summary box
    d_box = add_glass_card(slide14, Inches(0.8), Inches(6.85), Inches(11.733), Inches(0.45), bg_color=CARD_BG_ALT, border_color=CYAN)
    dtf = d_box.text_frame
    p = dtf.paragraphs[0]
    p.text = "🎬 Live Demo Focus: Demonstrate real-time code parsing, multi-agent reasoning, interactive dashboard metrics, and one-click PDF report export."
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 15: Conclusion
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_bg(slide15)
    add_header(slide15, "CONCLUSION", "Transforming Developer Productivity via Agentic AI", "Empowering software engineers with autonomous multi-agent intelligence.")
    add_anim_badge(slide15, "Final Keynote Outro & Logo Glow")

    # Left Summary Card
    l_card = add_glass_card(slide15, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1), bg_color=CARD_BG, border_color=BORDER_GLOW)
    ltf = l_card.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_top = ltf.margin_right = ltf.margin_bottom = Inches(0.2)

    p = ltf.paragraphs[0]
    p.text = "KEY TAKEAWAYS & IMPACT"
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE

    takeaways = [
        "Replaces fragmented static analysis tools with a collaborative Multi-Agent AI architecture.",
        "Significantly reduces debugging, security scanning, performance tuning, and code review cycles.",
        "Generates automated unit tests and explainable developer-friendly documentation.",
        "Empowers developers and software engineering teams to ship secure, high-quality code faster.",
        "Demonstrates the transformative potential of Agentic Coding in modern software development."
    ]

    for t in takeaways:
        pt = ltf.add_paragraph()
        pt.text = f"\n• {t}"
        pt.font.name = FONT_BODY
        pt.font.size = Pt(11)
        pt.font.color.rgb = TEXT_WHITE

    # Right Hero Closing Card
    r_card = add_glass_card(slide15, Inches(6.8), Inches(1.9), Inches(5.733), Inches(5.1), bg_color=CARD_BG_ALT, border_color=BORDER_PURPLE)
    rtf15 = r_card.text_frame
    rtf15.word_wrap = True
    rtf15.margin_left = rtf15.margin_top = rtf15.margin_right = rtf15.margin_bottom = Inches(0.25)

    p = rtf15.paragraphs[0]
    p.text = "CodeMedic AI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE

    p_sub = rtf15.add_paragraph()
    p_sub.text = "Fix. Explain. Optimize.\n"
    p_sub.font.name = FONT_TITLE
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = TEXT_WHITE

    p_pow = rtf15.add_paragraph()
    p_pow.text = "⚡ Powered by OpenAI Codex\n"
    p_pow.font.name = FONT_BODY
    p_pow.font.size = Pt(14)
    p_pow.font.bold = True
    p_pow.font.color.rgb = PURPLE

    p_dev = rtf15.add_paragraph()
    p_dev.text = "DEVELOPED BY: Umer Khan\nB.E. Artificial Intelligence & Data Science\n\nOPENAI CODEX HACKATHON 2026\nTrack: Agentic Coding"
    p_dev.font.name = FONT_BODY
    p_dev.font.size = Pt(12)
    p_dev.font.color.rgb = CYAN

    qp = rtf15.add_paragraph()
    qp.text = '\n\n“Empowering developers through Agentic Coding—building smarter software with autonomous AI agents.”'
    qp.font.name = FONT_TITLE
    qp.font.size = Pt(11.5)
    qp.font.bold = True
    qp.font.italic = True
    qp.font.color.rgb = TEXT_WHITE
    qp.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 16: Dedicated Project Links & Resources
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_bg(slide16)
    add_header(slide16, "PROJECT RESOURCES & LINKS", "CodeMedic AI Access & Repository Portal", "Explore the open-source GitHub codebase, live web application, technical documentation, and developer contact.")
    add_anim_badge(slide16, "Final Slide Pulse & Link Glow")

    # 4 Main Dedicated Link Cards (2x2 Grid)
    link_cards_data = [
        ("📂  GitHub Repository", "github.com/UmerKhan/CodeMedic-AI", SKY_BLUE, "Complete open-source repository including frontend, backend multi-agent services, Docker compose scripts, and ReportLab PDF generators."),
        ("🌐  Live Web Application", "codemedic-ai.vercel.app", GREEN, "Production application featuring Monaco Code Editor, SSE progress streaming, interactive health score dashboard, and multi-format report exports."),
        ("📄  API & Technical Docs", "docs.codemedic-ai.dev", PURPLE, "Comprehensive technical documentation detailing multi-agent schemas, FastAPI route endpoints, report mappers, and setup guides."),
        ("📧  Developer Contact", "umerkhan.dev@gmail.com", CYAN, "Umer Khan | B.E. Artificial Intelligence & Data Science\nOpenAI Codex Hackathon 2026 • Track: Agentic Coding")
    ]

    cw16 = Inches(5.7)
    ch16 = Inches(2.2)
    gx16 = Inches(0.33)
    gy16 = Inches(0.2)

    for i, (title, url, col_c, desc) in enumerate(link_cards_data):
        row = i // 2
        col = i % 2
        left = Inches(0.8) + col * (cw16 + gx16)
        top = Inches(1.9) + row * (ch16 + gy16)

        card = add_glass_card(slide16, left, top, cw16, ch16, bg_color=CARD_BG, border_color=BORDER_GLOW)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p_url = tf.add_paragraph()
        p_url.text = f"🔗 {url}"
        p_url.font.name = FONT_TITLE
        p_url.font.size = Pt(12)
        p_url.font.bold = True
        p_url.font.color.rgb = col_c

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_MUTED

    # Bottom Full-width Hero Banner with QR / Access Callout
    bot_card = add_glass_card(slide16, Inches(0.8), Inches(6.6), Inches(11.733), Inches(0.65), bg_color=CARD_BG_ALT, border_color=BORDER_PURPLE)
    btf16 = bot_card.text_frame
    p = btf16.paragraphs[0]
    p.text = "📱 Scan QR Code or Visit Links to Access Live Demo, Full Codebase & Download Technical Audit Reports."
    p.font.name = FONT_BODY
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.alignment = PP_ALIGN.CENTER

    out_file = "CodeMedic_AI_Hackathon_Presentation.pptx"
    try:
        prs.save(out_file)
        print(f"Presentation saved successfully: {out_file}")
    except PermissionError:
        out_file_v2 = "CodeMedic_AI_Hackathon_Presentation_V2.pptx"
        prs.save(out_file_v2)
        print(f"Original file locked by PowerPoint. Saved to: {out_file_v2}")

if __name__ == "__main__":
    create_presentation()
